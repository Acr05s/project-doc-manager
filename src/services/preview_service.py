"""文档预览服务 - 使用Python库实现Office文档的在线预览"""

import os
from io import BytesIO
import html
from docx import Document
import pandas as pd
from pptx import Presentation

class PreviewService:
    """文档预览服务类"""
    
    def docx_to_html(self, docx_path, page_number=None, total_pages=None):
        """将Word文档转换为HTML，支持分页预览"""
        try:
            doc = Document(docx_path)
            pages = []
            current_page = []
            page_size = 50  # 每页50个元素（段落+表格）
            element_count = 0
            
            # 处理标题和段落
            for paragraph in doc.paragraphs:
                if paragraph.style.name.startswith('Heading'):
                    level = int(paragraph.style.name[-1])
                    # 保留标题样式
                    style = ''
                    if hasattr(paragraph, 'alignment') and paragraph.alignment:
                        align_map = {
                            0: 'left',    # WD_ALIGN_PARAGRAPH.LEFT
                            1: 'center',  # WD_ALIGN_PARAGRAPH.CENTER
                            2: 'right',   # WD_ALIGN_PARAGRAPH.RIGHT
                            3: 'justify'  # WD_ALIGN_PARAGRAPH.JUSTIFY
                        }
                        align = align_map.get(paragraph.alignment, 'left')
                        style += f'text-align: {align};'
                    if style:
                        current_page.append(f'<h{level} style="{style}">{html.escape(paragraph.text)}</h{level}>')
                    else:
                        current_page.append(f'<h{level}>{html.escape(paragraph.text)}</h{level}>')
                else:
                    # 处理段落格式
                    style = ''
                    
                    # 处理对齐方式
                    if hasattr(paragraph, 'alignment') and paragraph.alignment:
                        align_map = {
                            0: 'left',    # WD_ALIGN_PARAGRAPH.LEFT
                            1: 'center',  # WD_ALIGN_PARAGRAPH.CENTER
                            2: 'right',   # WD_ALIGN_PARAGRAPH.RIGHT
                            3: 'justify'  # WD_ALIGN_PARAGRAPH.JUSTIFY
                        }
                        align = align_map.get(paragraph.alignment, 'left')
                        style += f'text-align: {align};'
                    
                    # 处理行距
                    if hasattr(paragraph, 'paragraph_format') and paragraph.paragraph_format:
                        if hasattr(paragraph.paragraph_format, 'space_after') and paragraph.paragraph_format.space_after:
                            space_after = paragraph.paragraph_format.space_after
                            style += f'margin-bottom: {space_after}pt;'
                        if hasattr(paragraph.paragraph_format, 'space_before') and paragraph.paragraph_format.space_before:
                            space_before = paragraph.paragraph_format.space_before
                            style += f'margin-top: {space_before}pt;'
                        # 处理缩进
                        if hasattr(paragraph.paragraph_format, 'left_indent') and paragraph.paragraph_format.left_indent:
                            left_indent = paragraph.paragraph_format.left_indent
                            style += f'margin-left: {left_indent}pt;'
                    
                    # 处理文本格式
                    runs_html = []
                    for run in paragraph.runs:
                        run_style = ''
                        
                        # 处理粗体
                        if run.bold:
                            run_style += 'font-weight: bold;'
                        
                        # 处理斜体
                        if run.italic:
                            run_style += 'font-style: italic;'
                        
                        # 处理下划线
                        if run.underline:
                            run_style += 'text-decoration: underline;'
                        
                        # 处理字体大小
                        if run.font.size:
                            font_size = run.font.size
                            run_style += f'font-size: {font_size.pt}pt;'
                        
                        # 处理字体颜色
                        if run.font.color and run.font.color.rgb:
                            rgb = run.font.color.rgb
                            hex_color = f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
                            run_style += f'color: {hex_color};'
                        
                        # 处理字体类型
                        if run.font.name:
                            run_style += f'font-family: {run.font.name}, sans-serif;'
                        
                        if run_style:
                            runs_html.append(f'<span style="{run_style}">{html.escape(run.text)}</span>')
                        else:
                            runs_html.append(html.escape(run.text))
                    
                    paragraph_html = ''.join(runs_html)
                    if style:
                        current_page.append(f'<p style="{style}">{paragraph_html}</p>')
                    else:
                        current_page.append(f'<p>{paragraph_html}</p>')
                
                element_count += 1
                if element_count >= page_size:
                    pages.append('\n'.join(current_page))
                    current_page = []
                    element_count = 0
            
            # 处理表格
            for table in doc.tables:
                table_html = []
                table_html.append('<table border="1" style="border-collapse: collapse; width: 100%; margin: 10px 0;">')
                for row in table.rows:
                    table_html.append('<tr>')
                    for cell in row.cells:
                        # 处理单元格内容
                        cell_content = []
                        for para in cell.paragraphs:
                            para_html = []
                            for run in para.runs:
                                run_style = ''
                                if run.bold:
                                    run_style += 'font-weight: bold;'
                                if run.italic:
                                    run_style += 'font-style: italic;'
                                if run.underline:
                                    run_style += 'text-decoration: underline;'
                                if run.font.size:
                                    font_size = run.font.size
                                    run_style += f'font-size: {font_size.pt}pt;'
                                if run.font.color and run.font.color.rgb:
                                    rgb = run.font.color.rgb
                                    hex_color = f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
                                    run_style += f'color: {hex_color};'
                                if run_style:
                                    para_html.append(f'<span style="{run_style}">{html.escape(run.text)}</span>')
                                else:
                                    para_html.append(html.escape(run.text))
                            cell_content.append(''.join(para_html))
                        cell_text = ''.join(cell_content)
                        table_html.append(f'<td style="padding: 5px;">{cell_text}</td>')
                    table_html.append('</tr>')
                table_html.append('</table>')
                
                current_page.append('\n'.join(table_html))
                element_count += 1
                if element_count >= page_size:
                    pages.append('\n'.join(current_page))
                    current_page = []
                    element_count = 0
            
            # 添加最后一页
            if current_page:
                pages.append('\n'.join(current_page))
            
            # 如果指定了页码，返回对应页
            if page_number is not None and 0 <= page_number < len(pages):
                return pages[page_number], len(pages)
            else:
                # 默认返回第一页
                return pages[0] if pages else '', len(pages)
        except Exception as e:
            return f'<p>预览失败: {str(e)}</p>', 1
    
    def excel_to_html(self, excel_path):
        """将Excel文档转换为HTML"""
        try:
            # 读取所有工作表
            xl_file = pd.ExcelFile(excel_path)
            html_content = []
            
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(xl_file, sheet_name=sheet_name)
                # 转换为HTML表格，保留更多格式
                sheet_html = f'<h2 style="margin-top: 20px; margin-bottom: 15px; padding-bottom: 5px; border-bottom: 1px solid #ddd;">{sheet_name}</h2>'
                sheet_html += df.to_html(
                    index=False, 
                    border=1, 
                    classes='excel-table',
                    na_rep='',
                    justify='left',
                    table_id=f'sheet-{sheet_name.replace(" ", "-")}',
                    render_links=True
                )
                html_content.append(sheet_html)
            
            return '\n'.join(html_content)
        except Exception as e:
            return f'<p>预览失败: {str(e)}</p>'
    
    def pptx_to_html(self, pptx_path):
        """将PPT文档转换为HTML"""
        try:
            prs = Presentation(pptx_path)
            html_content = []
            
            for i, slide in enumerate(prs.slides):
                html_content.append(f'<div style="margin-top: 30px; padding: 20px; border: 1px solid #ddd; border-radius: 8px; background: #f9f9f9;">')
                html_content.append(f'<h2 style="margin-top: 0; margin-bottom: 20px; padding-bottom: 10px; border-bottom: 1px solid #ddd;">幻灯片 {i+1}</h2>')
                
                # 处理形状
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        # 尝试获取形状的位置和大小信息
                        shape_style = ''
                        if hasattr(shape, 'left') and hasattr(shape, 'top'):
                            shape_style += f'position: relative; margin-bottom: 15px;'
                        
                        # 处理文本格式
                        text_frame = shape.text_frame
                        if text_frame:
                            paragraphs_html = []
                            for paragraph in text_frame.paragraphs:
                                para_style = ''
                                
                                # 处理段落对齐
                                if hasattr(paragraph, 'alignment') and paragraph.alignment:
                                    align_map = {
                                        0: 'left',    # PP_ALIGN.LEFT
                                        1: 'center',  # PP_ALIGN.CENTER
                                        2: 'right',   # PP_ALIGN.RIGHT
                                        3: 'justify'  # PP_ALIGN.JUSTIFY
                                    }
                                    align = align_map.get(paragraph.alignment, 'left')
                                    para_style += f'text-align: {align};'
                                
                                # 处理段落内容
                                runs_html = []
                                for run in paragraph.runs:
                                    run_style = ''
                                    
                                    # 处理粗体
                                    if run.bold:
                                        run_style += 'font-weight: bold;'
                                    
                                    # 处理斜体
                                    if run.italic:
                                        run_style += 'font-style: italic;'
                                    
                                    # 处理下划线
                                    if run.underline:
                                        run_style += 'text-decoration: underline;'
                                    
                                    # 处理字体大小
                                    if run.font.size:
                                        font_size = run.font.size
                                        run_style += f'font-size: {font_size.pt}pt;'
                                    
                                    # 处理字体颜色
                                    if run.font.color and run.font.color.rgb:
                                        rgb = run.font.color.rgb
                                        hex_color = f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
                                        run_style += f'color: {hex_color};'
                                    
                                    # 处理字体类型
                                    if run.font.name:
                                        run_style += f'font-family: {run.font.name}, sans-serif;'
                                    
                                    if run_style:
                                        runs_html.append(f'<span style="{run_style}">{html.escape(run.text)}</span>')
                                    else:
                                        runs_html.append(html.escape(run.text))
                                
                                paragraph_html = ''.join(runs_html)
                                if para_style:
                                    paragraphs_html.append(f'<p style="{para_style}">{paragraph_html}</p>')
                                else:
                                    paragraphs_html.append(f'<p>{paragraph_html}</p>')
                            
                            shape_html = ''.join(paragraphs_html)
                            if shape_style:
                                html_content.append(f'<div style="{shape_style}">{shape_html}</div>')
                            else:
                                html_content.append(shape_html)
                
                html_content.append('</div>')
            
            return '\n'.join(html_content)
        except Exception as e:
            return f'<p>预览失败: {str(e)}</p>'
    
    def get_preview_html(self, file_path, page_number=None):
        """根据文件类型获取预览HTML，支持分页预览"""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext in ['.docx']:
            content, total_pages = self.docx_to_html(file_path, page_number)
            return content, total_pages
        elif ext in ['.xlsx', '.xls']:
            return self.excel_to_html(file_path), 1
        elif ext in ['.pptx', '.ppt']:
            return self.pptx_to_html(file_path), 1
        else:
            return '<p>不支持的文件类型</p>', 1
    
    def get_full_preview(self, file_path, page_number=0):
        """获取完整的预览HTML页面，支持分页导航"""
        content, total_pages = self.get_preview_html(file_path, page_number)
        
        # 生成分页导航
        pagination_html = ''
        if total_pages > 1:
            pagination_html = f'''
            <div class="pagination" style="margin: 20px 0; text-align: center;">
                <button onclick="changePage({page_number - 1})" {'disabled' if page_number == 0 else ''} style="margin: 0 5px; padding: 5px 10px; border: 1px solid #ddd; border-radius: 4px; background: #f8f9fa; cursor: pointer;">上一页</button>
                <span style="margin: 0 10px;">第 {page_number + 1} 页，共 {total_pages} 页</span>
                <button onclick="changePage({page_number + 1})" {'disabled' if page_number >= total_pages - 1 else ''} style="margin: 0 5px; padding: 5px 10px; border: 1px solid #ddd; border-radius: 4px; background: #f8f9fa; cursor: pointer;">下一页</button>
            </div>
            '''
        
        full_html = f'''
        <!DOCTYPE html>
        <html>
        <head>
            <title>文档预览</title>
            <style>
                body {{
                    font-family: Arial, sans-serif;
                    margin: 20px;
                    line-height: 1.6;
                    background-color: white;
                    color: #333;
                }}
                h1, h2, h3, h4, h5, h6 {{
                    color: #333;
                    margin-top: 1.5em;
                    margin-bottom: 0.5em;
                }}
                p {{
                    margin-top: 0.5em;
                    margin-bottom: 1em;
                }}
                table {{
                    border-collapse: collapse;
                    width: 100%;
                    margin: 15px 0;
                    font-size: 14px;
                }}
                th, td {{
                    border: 1px solid #ddd;
                    padding: 8px 12px;
                    text-align: left;
                    vertical-align: top;
                }}
                th {{
                    background-color: #f2f2f2;
                    font-weight: bold;
                }}
                .excel-table {{
                    width: 100%;
                    font-size: 14px;
                    border-collapse: collapse;
                }}
                .excel-table th {{
                    background-color: #e8f4f8;
                    font-weight: bold;
                    text-align: center;
                }}
                .excel-table td {{
                    padding: 6px 10px;
                }}
                .excel-table tr:nth-child(even) {{
                    background-color: #f9f9f9;
                }}
                .preview-content {{
                    max-width: 900px;
                    margin: 0 auto;
                }}
                /* 响应式表格 */
                @media screen and (max-width: 768px) {{
                    table {{
                        font-size: 12px;
                    }}
                    th, td {{
                        padding: 4px 6px;
                    }}
                }}
            </style>
            <script>
                function changePage(page) {{
                    // 获取当前URL
                    const url = new URL(window.location.href);
                    // 设置page参数
                    url.searchParams.set('page', page);
                    // 刷新页面
                    window.location.href = url.toString();
                }}
            </script>
        </head>
        <body>
            <h1>文档预览</h1>
            {pagination_html}
            <div class="preview-content">
                {content}
            </div>
            {pagination_html}
        </body>
        </html>
        '''
        
        return full_html
